from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# פונקציה לקיבוע הטקסט לימין (RTL)
def set_rtl(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set('rtl', '1')

# פונקציית עזר להוספת תבליט מיושר לימין (הועברה לכאן כדי למנוע שגיאות)
def add_bullet(text_frame, text):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.RIGHT
    set_rtl(p)
    return p

prs = Presentation()

base_path = r"docs\presentation\output" 
os.makedirs(base_path, exist_ok=True)

# ==========================================
# שקף 1: שער ומטרת הפרויקט
# ==========================================
slide1 = prs.slides.add_slide(prs.slide_layouts[5])

title1 = slide1.shapes.title
title1.text = "פיתוח מערכת מבוססת FPGA להתראת נפילה לחולי דמנציה"
set_rtl(title1.text_frame.paragraphs[0])
title1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_desc = slide1.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.5))
p_desc = txBox_desc.text_frame.paragraphs[0]
p_desc.text = "קליטת תמונה בזמן אמת ממצלמת OV7670 והצגתה דרך כרטיס Artix-7."
p_desc.font.size = Pt(20)
p_desc.alignment = PP_ALIGN.CENTER
set_rtl(p_desc)

# הוספת שורה שנייה עם ירידת שורה מפורשת
p_desc2 = txBox_desc.text_frame.add_paragraph()
p_desc2.text = "על מנת לאפשר עיבוד תמונה לזיהויי עמידה ממושכת לחולי דמנציה."
p_desc2.font.size = Pt(17)
p_desc2.alignment = PP_ALIGN.CENTER
set_rtl(p_desc2)

# שליפת תמונת הבורד ישירות מהנתיב המלא בפרויקט
image_path1 = "docs/presentation/assets/artix_board.png"

if os.path.exists(image_path1):
    slide1.shapes.add_picture(image_path1, Inches(2.5), Inches(2.0), width=Inches(5))
else:
    fallback_box = slide1.shapes.add_textbox(Inches(2.5), Inches(3.0), Inches(5), Inches(1))
    fallback_p = fallback_box.text_frame.paragraphs[0]
    fallback_p.text = "[שגיאה: התמונה artix_board.png לא נמצאה]"
    fallback_p.font.size = Pt(14)
    fallback_p.font.color.rgb = RGBColor(255, 0, 0)
    fallback_p.alignment = PP_ALIGN.CENTER
    
txBox_details = slide1.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.5))
p_details = txBox_details.text_frame.paragraphs[0]
p_details.text = "מגיש: יוסי ברים | הנדסת חשמל שנה ד| המרכז האקדמי לב"
p_details.font.size = Pt(22)
p_details.font.bold = True
p_details.alignment = PP_ALIGN.CENTER
set_rtl(p_details)

# ==========================================
# שקף 2: ארכיטקטורת המערכת
# ==========================================
slide2 = prs.slides.add_slide(prs.slide_layouts[5])

title2 = slide2.shapes.title
title2.text = "ארכיטקטורת המערכת"
set_rtl(title2.text_frame.paragraphs[0])

txBox_text2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1))
tf_text2 = txBox_text2.text_frame

b1_s2 = tf_text2.add_paragraph()
b1_s2.text = "נתוני התמונה הגולמיים נקלטים מהמצלמה ונאגרים בזיכרון ה-BRAM."
b1_s2.font.size = Pt(18); b1_s2.alignment = PP_ALIGN.RIGHT; set_rtl(b1_s2)

b2_s2 = tf_text2.add_paragraph()
b2_s2.text = "בקר ה-VGA שולף את הנתונים ומעביר לממיר D to A לקבלת אות אנלוגי פיזי למסך."
b2_s2.font.size = Pt(18); b2_s2.alignment = PP_ALIGN.RIGHT; set_rtl(b2_s2)

top_box = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.8), Inches(2.4), Inches(2.2), Inches(1.8))
top_box.fill.solid(); top_box.fill.fore_color.rgb = RGBColor(0, 0, 0)
top_text = top_box.text_frame.paragraphs[0]
top_text.text = "מערכת התראת נפילה\nמבוסס FPGA"
top_text.font.color.rgb = RGBColor(255, 255, 255); top_text.alignment = PP_ALIGN.CENTER; set_rtl(top_text)

inputs_box = slide2.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(2.5), Inches(2))
inputs_box.text_frame.paragraphs[0].text = "Inputs:\nclk, reset\nov7670_vsync, href, pclk\nov7670_data[7:0]\nbtn[1:0], sw[1:0]\nscl, sda"
inputs_box.text_frame.paragraphs[0].font.size = Pt(13); inputs_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.1), Inches(3.1), Inches(0.6), Inches(0.3))

outputs_box = slide2.shapes.add_textbox(Inches(6.2), Inches(2.1), Inches(3.5), Inches(2))
outputs_box.text_frame.paragraphs[0].text = "Outputs:\nVGA_HS_O, VGA_VS_O\nVGA_R[3:0], G[3:0], B[3:0]\nov7670_xclk, pwdn, reset\nled[3:0]"
outputs_box.text_frame.paragraphs[0].font.size = Pt(13)
slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.1), Inches(3.1), Inches(0.6), Inches(0.3))

blocks = ["ov7670_capture", "frame_buffer", "vga_controller", "D to A Converter"]
for i, b_text in enumerate(blocks):
    b = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6 + i*2.2), Inches(5.2), Inches(1.7), Inches(0.8))
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(79, 129, 189)
    bp = b.text_frame.paragraphs[0]
    bp.text = b_text; bp.font.size = Pt(14); bp.alignment = PP_ALIGN.CENTER
    if i < 3:
        slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.3 + i*2.2), Inches(5.5), Inches(0.5), Inches(0.2))


# ==========================================
# שקף תוספת: רקע היסטורי וקונספט תותח האלקטרונים (CRT)
# ==========================================
slide_history = prs.slides.add_slide(prs.slide_layouts[5])
title_history = slide_history.shapes.title
title_history.text = "מבוא ל-VGA: קונספט תותח האלקטרונים (CRT)"
set_rtl(title_history.text_frame.paragraphs[0])
title_history.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_history = slide_history.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
tf_history = txBox_history.text_frame

add_bullet(tf_history, "• מקור הפרוטוקול: טכנולוגיית ה-VGA פותחה במקור עבור מסכי שפופרת קרן קתודית (CRT) המבוססים על פיזיקה אנלוגית.")
add_bullet(tf_history, "• תותח אלקטרונים: המסך פועל באמצעות קרן פיזית הנורית על גבי מסך מצופה זרחן. עוצמת הזרם קובעת את עוצמת ההארה של כל פיקסל.")
add_bullet(tf_history, "• המורשת האנלוגית: למרות שכיום המסכים מודרניים, הפרוטוקול מחייב שידור אותות השהייה שיאפשרו לקרן הפיזית זמן תנועה וחזרה.")

# תמונת תותח האלקטרונים
image_path_history = "docs/presentation/assets/vga_analog_cathode_ray_concept.png" 
if os.path.exists(image_path_history):
    slide_history.shapes.add_picture(image_path_history, Inches(2.5), Inches(3.2), width=Inches(5.0))

# ==========================================
# שקף 3: צעד 1 בסיפור - הפיזיקה של הסריקה הקווית (Raster Scan)
# ==========================================
slide3 = prs.slides.add_slide(prs.slide_layouts[5])
title3 = slide3.shapes.title
title3.text = "מבוא ל-VGA: סריקה קווית (Raster Scan)"
set_rtl(title3.text_frame.paragraphs[0])
title3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
tf_text3 = txBox_text3.text_frame

add_bullet(tf_text3, "• הציור על המסך מבוסס על קרן שסורקת את הפיקסלים משמאל לימין (שורה) ומלמעלה למטה (פריים).")
add_bullet(tf_text3, "• תנועת הקרן חייבת להיות רציפה, מה שיוצר מסלול בצורת 'זיגזג'.")
add_bullet(tf_text3, "• בכל פעם שהקרן מגיעה לקצה (אופקי או אנכי), נדרש זמן כדי להחזיר אותה לנקודת ההתחלה (Retrace).")

# תמונת הסריקה הקווית הממחישה את הזיגזג וה-Retrace
image_path3 = "docs/presentation/assets/vga_raster_scan_retrace_concept.png" 
if os.path.exists(image_path3):
    slide3.shapes.add_picture(image_path3, Inches(2.5), Inches(3.0), width=Inches(4.5))


# ==========================================
# שקף תוספת (שקף 5 החדש): תזמוני VGA - הפיזיקה מאחורי המספרים
# ==========================================
slide_table = prs.slides.add_slide(prs.slide_layouts[5])
title_table = slide_table.shapes.title
title_table.text = "תזמוני VGA: הפיזיקה מאחורי המספרים"
set_rtl(title_table.text_frame.paragraphs[0])
title_table.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# יצירת טבלה 5 שורות, 4 עמודות
x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.0)
table_shape = slide_table.shapes.add_table(5, 4, x, y, cx, cy)
table = table_shape.table

# הגדרת רוחב העמודות
table.columns[0].width = Inches(1.5) # פרמטר
table.columns[1].width = Inches(1.5) # שעונים
table.columns[2].width = Inches(3.0) # משמעות היסטורית
table.columns[3].width = Inches(3.0) # משמעות מעשית (FPGA)

# כותרות הטבלה
headers = ["פרמטר", "מחזורי שעון", "הצורך הפיזיקלי המקורי (CRT)", "המימוש שלנו (RTL)"]
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(16)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_rtl(cell.text_frame.paragraphs[0])

# תוכן השורות
row_data = [
    ["Active Display", "640", "הקרן מציירת את הפיקסלים הפיזיים על המסך משמאל לימין.", "שליפת נתונים מה-BRAM ושידור צבע (RGB) פעיל."],
    ["Front Porch", "16", "מרווח זמן לקרן 'להירגע' בסוף השורה לפני שחוזרת אחורה.", "כיבוי מיידי של ה-RGB ל-'0000' למניעת מריחות."],
    ["Sync Pulse", "96", "מכת מתח (טריגר) שמאלצת את הקרן לקפוץ חזרה שמאלה.", "הורדת האות הפיזי HSYNC ל-'0' למשך 96 שעונים."],
    ["Back Porch", "48", "המתנה לייצוב הקרן בתחילת השורה החדשה בצד שמאל.", "המשך השהיית ה-RGB על '0000' עד הגעה לפיקסל הראשון."]
]

for row_idx, row in enumerate(row_data):
    for col_idx, text in enumerate(row):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = text
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        if col_idx < 2:
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        else:
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        set_rtl(cell.text_frame.paragraphs[0])


# ==========================================
# שקף 5: תזמונים והחשכה בחומרה (VGA Controller)
# ==========================================
slide5 = prs.slides.add_slide(prs.slide_layouts[5])
title5 = slide5.shapes.title
title5.text = "תזמונים והחשכה בחומרה (VGA Controller)"
set_rtl(title5.text_frame.paragraphs[0])
title5.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

tx_c5 = slide5.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(1.5))
tf_c5 = tx_c5.text_frame

add_bullet(tf_c5, "• מחזור שלם (800 פיקסלים): חיבור הערכים (640+16+96+48) דורש מהמונה האופקי (h_count) לספור מ-0 ועד 799.")
add_bullet(tf_c5, "• שעון פיקסל 25MHz: כל 'פיקסל' שקול למחזור שעון אחד של 40ns. סנכרון זה מבטיח קצב ריענון תקני של 60Hz.")
add_bullet(tf_c5, "• אכיפת שחור (Blanking): מחוץ לתחום ה-640, הלוגיקה שלנו מאלצת \"0000\" ב-RGB כדי למנוע 'מריחות' צבע בעת תנועת הקרן.")

# 1. דיאגרמת הבלוקים והמונים (RTL Schematic)
img1_path5 = "docs/presentation/assets/vga_counters_rtl_schematic.png"
if os.path.exists(img1_path5):
    slide5.shapes.add_picture(img1_path5, Inches(0.5), Inches(3.2), width=Inches(4.5))

# 2. גרף התזמונים והטבלה שממנה נגזרים הערכים
img2_path5 = "docs/presentation/assets/vga_timing_waveform_and_table.png"
if os.path.exists(img2_path5):
    slide5.shapes.add_picture(img2_path5, Inches(5.2), Inches(2.8), width=Inches(4.0))

# ==========================================
# שקף 6: מסלול הנתונים - חישוב כתובת ושליפה מה-BRAM
# ==========================================
slide6 = prs.slides.add_slide(prs.slide_layouts[5])
title6 = slide6.shapes.title
title6.text = "מסלול הנתונים: חישוב כתובת הפיקסל ושליפה"
set_rtl(title6.text_frame.paragraphs[0])
title6.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
tf_text6 = txBox_text6.text_frame

add_bullet(tf_text6, "• מיקום הקרן (X,Y): המונים hsync_reg ו-vsync_reg מפיקים את הקואורדינטות הדינמיות display_x ו-display_y בתוך האזור הפעיל.")
add_bullet(tf_text6, "• תרגום לכתובת ליניארית: הקואורדינטה מתורגמת מתמטית לכתובת הזיכרון bram_address_next ונשלחת החוצה דרך האות addrb.")
add_bullet(tf_text6, "• שליפה מה-BRAM: רכיב ה-frame_buffer מקבל את הכתובת ומשחרר את 12 הביטים של הפיקסל היישר אל האות doutb של בקר ה-VGA.")

img_path6 = "docs/presentation/assets/vga_address_bram_flow.png" 
if os.path.exists(img_path6):
    slide6.shapes.add_picture(img_path6, Inches(2.5), Inches(3.2), width=Inches(5.0))

# ==========================================
# שקף 7: פענוח הצבע ובקרת הוידאו
# ==========================================
slide7 = prs.slides.add_slide(prs.slide_layouts[5])
title7 = slide7.shapes.title
title7.text = "עיבוד הפיקסל: ניתוב ל-RGB ומנגנון ה-Blanking"
set_rtl(title7.text_frame.paragraphs[0])
title7.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text7 = slide7.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
tf_text7 = txBox_text7.text_frame

add_bullet(tf_text7, "• רישום ופיצול: הנתונים מ-doutb ננעלים באוגר pxl_data_reg ומפוצלים לשגרירים: vga_r_int (ביטים 11:8), vga_g_int (7:4) ו-vga_b_int (3:0).")
add_bullet(tf_text7, "• שומר הסף: האות in_display_area_delayed משמש כדגל המזהה האם הקרן מצוירת כעת על המסך או נמצאת מחוץ לתצוגה.")
add_bullet(tf_text7, "• סינון דיגיטלי (Mux): כאשר הדגל ב-'1', ערכי ה-RGB מועברים ליציאות VGA_R/G/B. בזמני ה-Porch וה-Sync, היציאות נאלצות ל-\"0000\" (שחור מוחלט).")

img_path7 = "docs/presentation/assets/vga_rgb_video_on_logic.png" 
if os.path.exists(img_path7):
    slide7.shapes.add_picture(img_path7, Inches(2.5), Inches(3.2), width=Inches(5.0))

# ==========================================
# שקף 8: המעבר לעולם האנלוגי ומחבר ה-VGA (DAC)
# ==========================================
slide8 = prs.slides.add_slide(prs.slide_layouts[5])
title8 = slide8.shapes.title
title8.text = "המרת DAC אנלוגית ומחבר ה-VGA"
set_rtl(title8.text_frame.paragraphs[0])
title8.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text8 = slide8.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(1.0))
tf_text8 = txBox_text8.text_frame

add_bullet(tf_text8, "• סולם נגדים (Resistor Network): המרת 12 ביטי ה-RGB למתח אנלוגי רציף באמצעות משקולות נגדים (510Ω עד 4KΩ).")
add_bullet(tf_text8, "• אותות סנכרון (B11 / B12): פיני היציאה מוגנים באמצעות נגדי 100Ω ומזינים ישירות את מחבר ה-DB15 של המסך.")

img_path8_1 = "docs/presentation/assets/vga_dac_resistor_network.png" 
if os.path.exists(img_path8_1):
    slide8.shapes.add_picture(img_path8_1, Inches(0.5), Inches(2.3), width=Inches(4.4))

img_path8_2 = "docs/presentation/assets/vga_db15_connector_pins.png" 
if os.path.exists(img_path8_2):
    slide8.shapes.add_picture(img_path8_2, Inches(5.1), Inches(2.5), width=Inches(4.4))

# ==========================================
# שקף 9: עומק הנדסי - הפיזיקה של מחלק המתח (MSB/LSB)
# ==========================================
slide9 = prs.slides.add_slide(prs.slide_layouts[5])
title9 = slide9.shapes.title
title9.text = "  מחלק המתח ומשקלי הביטים"
set_rtl(title9.text_frame.paragraphs[0])
title9.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text9 = slide9.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
tf_text9 = txBox_text9.text_frame

add_bullet(tf_text9, "• משקל פיזי לכל ביט (MSB לעומת LSB): הביט המשמעותי ביותר (MSB) מחובר לנגד הקטן ביותר, ולכן מזזרים את הזרם הגדול ביותר למעגל ומשפיע משמעותית על עוצמת הצבע.")
add_bullet(tf_text9, "• כוונון עדין: לעומתו, הביט הפחות משמעותי (LSB) מחובר לנגד הגדול ביותר בסולם. תרומתו לזרם הכולל היא מזערית ונועדה לכוונון עדין של הגוון.")
add_bullet(tf_text9, "• סכימה אנלוגית על המסך: כל הזרמים מהביטים הפעילים ('1') מתחברים וזורמים יחד דרך נגד הסיומת של המסך (75Ω) אל האדמה. כך נוצר מחלק מתח דינמי המפיק 0V עד 0.7V.")

img_path9 = "docs/presentation/assets/vga_voltage_divider_dac.png" 
if os.path.exists(img_path9):
    slide9.shapes.add_picture(img_path9, Inches(1.5), Inches(3.2), width=Inches(7.0))


# ==========================================
# שקף תוספת: סיכום המבוא - לקראת ארכיטקטורת החומרה
# ==========================================
slide10 = prs.slides.add_slide(prs.slide_layouts[5])
title10 = slide10.shapes.title
title10.text = "סיכום תיאורטי: לקראת ארכיטקטורת החומרה"
set_rtl(title10.text_frame.paragraphs[0])
title10.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text10 = slide10.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(1.5))
tf_text10 = txBox_text10.text_frame

add_bullet(tf_text10, "• האתגר הפיזיקלי: ראינו שהמסך החיצוני דורש תזמונים מדויקים ונוקשים (Sync, Blanking) והמרת צבע רציפה (DAC) כדי לפעול כראוי.")
add_bullet(tf_text10, "• המעבר פנימה (FPGA): כדי לייצר את האותות הללו בזמן אמת ובאמינות מוחלטת, אנו נכנסים אל תוך הלוגיקה הדיגיטלית (RTL) על הכרטיס.")
add_bullet(tf_text10, "• תפקיד ה-VGA Controller: זהו רכיב הגישור הסופי במערכת. הוא שואב נתונים מהזיכרון הפנימי ומתרגם אותם לאותות הפיזיים שהמסך דורש.")
add_bullet(tf_text10, "• התחנה הבאה: נמפה את האקו-סיסטם על הכרטיס (שעונים, מצלמה, חוצץ זיכרון) ולאחר מכן נצלול פנימה אל הלוגיקה של בקר ה-VGA עצמו.")

# הנתיב המעודכן לתמונת הארכיטקטורה
img_path10 = "docs/presentation/assets/vga_fpga_top_architecture.png"

if os.path.exists(img_path10):
    slide10.shapes.add_picture(img_path10, Inches(2.5), Inches(3.6), width=Inches(5.0))
else:
    err_box = slide10.shapes.add_textbox(Inches(2.5), Inches(4.0), Inches(5), Inches(1))
    err_p = err_box.text_frame.paragraphs[0]
    err_p.text = f"[שגיאה: התמונה לא נמצאה בנתיב: {img_path10}]"
    err_p.font.size = Pt(14)
    err_p.font.color.rgb = RGBColor(255, 0, 0)
    err_p.alignment = PP_ALIGN.CENTER
    
# ==========================================
# שקף 10: מבט-על - זרימת נתונים ותחומי שעון (RTL)
# ==========================================
slide10 = prs.slides.add_slide(prs.slide_layouts[5])
title10 = slide10.shapes.title
title10.text = "מבט-על: זרימת נתונים ותחומי שעון (RTL)"
set_rtl(title10.text_frame.paragraphs[0])
title10.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text10 = slide10.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(1.5))
tf_text10 = txBox_text10.text_frame

add_bullet(tf_text10, "• מחולל שעונים (clk_generator): רכיב PLL המקבל שעון מערכת ומפיק שעון 25MHz נפרד ל-VGA ושעון ייעודי למצלמה.")
add_bullet(tf_text10, "• גישור חציית שעונים (CDC): זיכרון ה-BRAM (בתצורת Dual-Port) פועל כחוצץ בטוח בין קצב הכתיבה לקצב הקריאה.")
add_bullet(tf_text10, "• תחום אדום (Camera Domain): קליטת הפיקסלים מהמצלמה (ov7670_capture) ודחיפתם לפורט A של הזיכרון (clka).")
add_bullet(tf_text10, "• תחום כחול (VGA Domain): המיקוד שלנו – שליפת הפיקסלים מפורט B (clkb) על ידי בקר ה-VGA וייצור אותות התצוגה.")

# תמונת ה-RTL
img_path10 = "docs/presentation/assets/rtl_direct_vga_path.png"
if os.path.exists(img_path10):
    # מיקום התמונה (ממורכזת)
    slide10.shapes.add_picture(img_path10, Inches(0.5), Inches(3.0), width=Inches(9.0))
    
    # יצירת מסגרת אדומה (Camera Domain) - צד שמאל ועד אמצע הזיכרון
    red_box = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.2), Inches(4.0), Inches(3.8))
    red_box.fill.background() # הופך את תוכן הריבוע לשקוף
    red_box.line.color.rgb = RGBColor(255, 0, 0) # צבע אדום
    red_box.line.width = Pt(3)
    
    # יצירת מסגרת כחולה (VGA Domain - התחום שלך) - מאמצע הזיכרון ועד ימין
    blue_box = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), Inches(3.2), Inches(4.0), Inches(3.8))
    blue_box.fill.background() # הופך את תוכן הריבוע לשקוף
    blue_box.line.color.rgb = RGBColor(0, 112, 192) # צבע כחול
    blue_box.line.width = Pt(4) # קצת יותר עבה כדי להדגיש שזה התחום שלך


# ==========================================
# שקף 11: צלילה לתכן - ניהול שעונים (Clock Generator)
# ==========================================
slide11 = prs.slides.add_slide(prs.slide_layouts[5])
title11 = slide11.shapes.title
title11.text = "ניהול שעונים במערכת (Clock Generator)"
set_rtl(title11.text_frame.paragraphs[0])
title11.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text11 = slide11.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(1.5))
tf_text11 = txBox_text11.text_frame

# שימוש ממוקד במספרים ושמות האותות מתוך הלוגיקה
add_bullet(tf_text11, "• מקור שעון (System Clock): כניסת clk מקבלת את שעון הלוח המקורי (100MHz) ומזינה את ה-PLL (clk_wiz_0).")
add_bullet(tf_text11, "• שעון תצוגה (vga_pll): ה-PLL מפיק תדר של 25MHz המזין את בקר ה-VGA (pxl_clk) ואת פורט הקריאה (clkb) ב-BRAM.")
add_bullet(tf_text11, "• שעון מצלמה (xclk_pll): הפקת שעון ייעודי של 24MHz (xclk_ov7670) המנותב החוצה לסנכרון רכיב המצלמה.")

# הכנה לתמונה מתאימה (סניפט קוד או סכמת בלוק של clk_wiz_0)
img_path11 = "docs/presentation/assets/clk_wiz_instantiation.png"
if os.path.exists(img_path11):
    slide11.shapes.add_picture(img_path11, Inches(2.5), Inches(3.2), width=Inches(5.0))
    

# ==========================================
# שקף 12: חוצץ התמונה - זיכרון ה-BRAM (Frame Buffer)
# ==========================================
slide12 = prs.slides.add_slide(prs.slide_layouts[5])
title12 = slide12.shapes.title
title12.text = "זיכרון ה-BRAM: חוצץ התמונה (Frame Buffer)"
set_rtl(title12.text_frame.paragraphs[0])
title12.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text12 = slide12.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(1.5))
tf_text12 = txBox_text12.text_frame

add_bullet(tf_text12, "• תצורת Dual-Port: הפרדה מוחלטת לחומרת כתיבה (A) וקריאה (B). זהו הפתרון שלנו ל-CDC (חציית תחומי שעון) המונע התנגשויות נתונים.")
add_bullet(tf_text12, "• שעונים א-סינכרוניים: הכתיבה (clka) מתבצעת בתדר המערכת (100MHz), בעוד הקריאה למסך (clkb) מונעת בלעדית על ידי שעון ה-VGA שלנו (25MHz).")
add_bullet(tf_text12, "• ממדי הזיכרון (307,200x12): הזיכרון מוקצה במדויק להכיל 307,200 פיקסלים (רזולוציה של 640x480).")
add_bullet(tf_text12, "• תפוקת הנתונים: האות doutb משחרר 12 ביט של צבע (RGB) בכל מחזור שעון (25MHz) ישירות לבקר ה-VGA.")

# הכנה לתמונת ה-IP Symbol של רכיב ה-BRAM (blk_mem_gen_0)
img_path12 = "docs/presentation/assets/bram_ip_symbol.png"
if os.path.exists(img_path12):
    slide12.shapes.add_picture(img_path12, Inches(2.5), Inches(3.5), width=Inches(5.0))
    

# ==========================================
# שקף 13: הצד האדום - לכידת התמונה (Camera to BRAM)
# ==========================================
slide13 = prs.slides.add_slide(prs.slide_layouts[5])
title13 = slide13.shapes.title
title13.text = "ממשק המצלמה: מהעולם הפיזי אל הזיכרון"
set_rtl(title13.text_frame.paragraphs[0])
title13.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text13 = slide13.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(1.5))
tf_text13 = txBox_text13.text_frame

add_bullet(tf_text13, "• הגשר הפיזי (XCLK מול PCLK): ה-FPGA מפיק 24MHz (XCLK) החוצה להפעלת המצלמה. בתגובה, המצלמה משדרת חזרה פיקסלים יחד עם שעון סנכרון עצמאי משלה (PCLK).")
add_bullet(tf_text13, "• לכידת הנתונים: רכיב ה-ov7670_capture רץ בתדר המערכת (100MHz) המאפשר דגימה מהירה ובטוחה של ה-PCLK ואסיפת 12 ביט של צבע.")
add_bullet(tf_text13, "• כתיבה לזיכרון: הבלוק מייצר את הכתובת (addra), המידע (dina), ודגל הכתיבה (wea) ודוחף אותם לפורט A של ה-BRAM.")

# הכנה לתמונה הממחישה את הלכידה
img_path13 = "docs/presentation/assets/camera_capture_flow.png.jpeg"
 
if os.path.exists(img_path13):
    slide13.shapes.add_picture(img_path13, Inches(2.5), Inches(3.5), width=Inches(5.0))


# ==========================================
# שקף 15: ממשק הזיכרון - לוגיקת שליפת הנתונים
# ==========================================
slide15 = prs.slides.add_slide(prs.slide_layouts[5])
title15 = slide15.shapes.title
title15.text = "ממשק הזיכרון (BRAM): ממדים וסנכרון נתונים"
set_rtl(title15.text_frame.paragraphs[0])
title15.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text15 = slide15.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(2.0))
tf_text15 = txBox_text15.text_frame

add_bullet(tf_text15, "• תרגום קואורדינטות: הכתובת לזיכרון (addrb באורך 19 ביט) מחושבת מתמטית לפי המיקום באזור הפעיל: Y * 640 + X.")
add_bullet(tf_text15, "• מנגנון ה-Pipeline (עיכוב מכוון): הקריאה מהזיכרון א-סינכרונית. כדי להבטיח יציבות, הקוד נועל את הנתון הנכנס (doutb) לתוך אוגר פנימי (pxl_data_reg) בעליית השעון הבאה.")
add_bullet(tf_text15, "• סנכרון שומר הסף: כדי לפצות על השהיית הזיכרון, דגל אזור התצוגה (in_display_area) מושהה גם הוא במחזור שעון אחד (ל-in_display_area_delayed).")
add_bullet(tf_text15, "• ניתוב דיגיטלי (Mux): רק כאשר הדגל המושהה דולק, 12 הביטים מפוצלים ומנותבים ליציאות VGA_R, VGA_G ו-VGA_B (4 ביטים לצבע).")

# אפשר לשים כאן סניפט (צילום מסך של קטע הקוד של ה-PROCESS שמחשב את bram_address_next)
img_path15 = "docs/presentation/assets/vga_address_logic_code.png" # הכן צילום מסך של הקוד
if os.path.exists(img_path15):
    slide15.shapes.add_picture(img_path15, Inches(2.5), Inches(4.0), width=Inches(5.0))


# ==========================================
# שקף 16: סימולציית מסלול הנתונים (Data Path)
# ==========================================
slide16 = prs.slides.add_slide(prs.slide_layouts[5])
title16 = slide16.shapes.title
title16.text = "אימות מסלול הנתונים: VGA Controller ↔ BRAM"
set_rtl(title16.text_frame.paragraphs[0])
title16.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text16 = slide16.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(1.5))
tf_text16 = txBox_text16.text_frame

add_bullet(tf_text16, "• מטרת הסימולציה: הוכחת לחיצת היד (Handshake) שבין בקשת הכתובת (addrb) לקבלת נתוני הפיקסל (doutb).")
add_bullet(tf_text16, "• דגימה בזמן אמת: ניתן לראות בבירור כיצד הכתובת עולה, המידע מופיע מיידית ב-doutb (קריאה א-סינכרונית), וננעל באוגר pxl_data_reg בעליית השעון הבאה.")
add_bullet(tf_text16, "• מניעת זבל ויזואלי: הסימולציה מוכיחה שבמעבר ל-Blanking (סוף שורה), יציאות ה-RGB נחתכות ל-0 מיידית, למרות שהזיכרון עדיין פולט מידע ישן.")

# כאן תיכנס התמונה של ה-Waveform מ-ModelSim שאותה נייצר כעת
img_path16 = "docs/presentation/assets/modelsim_data_path_flow.png" 
if os.path.exists(img_path16):
    slide16.shapes.add_picture(img_path16, Inches(1.0), Inches(3.6), width=Inches(8.0))
else:
    err_box = slide16.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(5), Inches(1))
    err_p = err_box.text_frame.paragraphs[0]
    err_p.text = f"[שומר מקום: כאן תיכנס תמונת ה-Waveform של זרימת הנתונים]"
    err_p.font.size = Pt(14)
    err_p.font.color.rgb = RGBColor(0, 112, 192)
    err_p.alignment = PP_ALIGN.CENTER


# ==========================================
# שקף 17: מבוא לוריפיקציה ומחזור שורה שלם (Ts)
# ==========================================
slide17 = prs.slides.add_slide(prs.slide_layouts[5])
title17 = slide17.shapes.title
title17.text = "מבוא לאימות התקן: מחזור שורה שלם (Ts)"
set_rtl(title17.text_frame.paragraphs[0])
title17.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text17 = slide17.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(1.5))
tf_text17 = txBox_text17.text_frame

add_bullet(tf_text17, "• שלב ב' - מהדיגיטל לפיזיקה: לאחר שהוכחנו את אמינות שליפת הנתונים מהזיכרון, כעת נוכיח שהחומרה מייצרת את התזמונים הפיזיים הנדרשים בתקן ה-VGA.")
add_bullet(tf_text17, "• מפת הדרכים (Macro to Micro): נתחיל ממבט-על של שורת מסך שלמה, ונצלול בהדרגה לניתוח האזורים הפעילים, זמני ההמתנה (Porches) ודופק הסנכרון.")
add_bullet(tf_text17, "• הוכחת מחזור השורה: התקן דורש 800 מחזורי שעון. מונה ה-hsync_reg סופר 0-799 במדויק, ומדידת הסמנים מוכיחה דלתא של 32us.")

img_path17 = "docs/presentation/assets/modelsim_full_line_ts.png" 
if os.path.exists(img_path17):
    slide17.shapes.add_picture(img_path17, Inches(1.0), Inches(3.6), width=Inches(8.0))
else:
    err_box = slide17.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(5), Inches(1))
    err_p = err_box.text_frame.paragraphs[0]; err_p.text = "[שומר מקום: modelsim_full_line_ts.png]"
    err_p.font.size = Pt(14); err_p.font.color.rgb = RGBColor(0, 112, 192); err_p.alignment = PP_ALIGN.CENTER

# ==========================================
# שקף 18: האזור הפעיל – זמן תצוגה
# ==========================================
slide18 = prs.slides.add_slide(prs.slide_layouts[5])
title18 = slide18.shapes.title
title18.text = "ליבת התצוגה: אימות האזור הפעיל (Tdisp)"
set_rtl(title18.text_frame.paragraphs[0])
title18.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text18 = slide18.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(1.5))
tf_text18 = txBox_text18.text_frame

add_bullet(tf_text18, "• תזמון הנתונים: מתוך כלל מחזור השורה, מוקצים בדיוק 640 מחזורי שעון לציור הפיקסלים (25.6us).")
add_bullet(tf_text18, "• דגל הבקרה: האות in_display_area עולה ל-'1' ומאפשר את שליפת הנתונים רק בטווח המורשה.")
add_bullet(tf_text18, "• מונה הפיקסלים הפעיל: המונה display_x סופר במדויק מ-0 ועד 639. לאורך כל החלון הזה, יציאות ה-RGB פולטות נתונים אקטיביים למסך.")

img_path18 = "docs/presentation/assets/modelsim_active_display_tdisp.png" 
if os.path.exists(img_path18):
    slide18.shapes.add_picture(img_path18, Inches(1.0), Inches(3.6), width=Inches(8.0))
else:
    err_box = slide18.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(5), Inches(1))
    err_p = err_box.text_frame.paragraphs[0]; err_p.text = "[שומר מקום: modelsim_active_display_tdisp.png]"
    err_p.font.size = Pt(14); err_p.font.color.rgb = RGBColor(0, 112, 192); err_p.alignment = PP_ALIGN.CENTER

# ==========================================
# שקף 19: השוליים ודחיקת האותות (Blanking)
# ==========================================
slide19 = prs.slides.add_slide(prs.slide_layouts[5])
title19 = slide19.shapes.title
title19.text = "אכיפת Blanking: שוליים (Porches)"
set_rtl(title19.text_frame.paragraphs[0])
title19.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text19 = slide19.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(1.5))
tf_text19 = txBox_text19.text_frame

add_bullet(tf_text19, "• משמעות ה'זמנים המתים': ה-Front Porch (16 שעונים) וה-Back Porch (48 שעונים) נועדו לייצוב תנועת הקרן.")
add_bullet(tf_text19, "• אכיפה לוגית: בסימולציה ניתן לראות כי ברגע שהקרן חורגת מהאזור הפעיל (display_x מסתיים), המערכת מאלצת את כל ערוצי ה-RGB ל-'0000'.")
add_bullet(tf_text19, "• התוצאה: שחור מוחלט על המסך באזורי המעבר, מה שמונע 'מריחות' צבע ומאפשר למסך להסתנכרן בצורה חלקה.")

img_path19 = "docs/presentation/assets/modelsim_porches_blanking.png" 
if os.path.exists(img_path19):
    slide19.shapes.add_picture(img_path19, Inches(1.0), Inches(3.6), width=Inches(8.0))
else:
    err_box = slide19.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(5), Inches(1))
    err_p = err_box.text_frame.paragraphs[0]; err_p.text = "[שומר מקום: modelsim_porches_blanking.png]"
    err_p.font.size = Pt(14); err_p.font.color.rgb = RGBColor(0, 112, 192); err_p.alignment = PP_ALIGN.CENTER


# ==========================================
# שקף 19.5 (תוספת): המרפסת האחורית (Back Porch) וחזרת הצבע
# ==========================================
slide_bp = prs.slides.add_slide(prs.slide_layouts[5])
title_bp = slide_bp.shapes.title
title_bp.text = "אכיפת Blanking: המרפסת האחורית (Back Porch)"
set_rtl(title_bp.text_frame.paragraphs[0])
title_bp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text_bp = slide_bp.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(1.5))
tf_text_bp = txBox_text_bp.text_frame

add_bullet(tf_text_bp, "• זמן המתנה (Back Porch): לאחר סיום דופק הסנכרון, המערכת ממתינה 48 מחזורי שעון (1.92us) לייצוב הקרן בתחילת השורה.")
add_bullet(tf_text_bp, "• המשך החשכה: בדומה למרפסת הקדמית, ה-RGB מוחזק במכוון על '0000' כדי למנוע זליגות צבע.")
add_bullet(tf_text_bp, "• חזרת הצבעים: ברגע שהאות in_display_area עולה ל-'1', רואים בבירור את נתוני ה-RGB האמיתיים מהזיכרון פורצים החוצה ומציירים את הפיקסל הראשון.")

img_path_bp = "docs/presentation/assets/modelsim_back_porch.png" 
if os.path.exists(img_path_bp):
    slide_bp.shapes.add_picture(img_path_bp, Inches(1.0), Inches(3.6), width=Inches(8.0))
else:
    err_box = slide_bp.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(5), Inches(1))
    err_p = err_box.text_frame.paragraphs[0]; err_p.text = "[שומר מקום: modelsim_back_porch.png]"
    err_p.font.size = Pt(14); err_p.font.color.rgb = RGBColor(0, 112, 192); err_p.alignment = PP_ALIGN.CENTER
    
    
# ==========================================
# שקף 20: דופק הסנכרון והוכחה אוטומטית (Tpw)
# ==========================================
slide20 = prs.slides.add_slide(prs.slide_layouts[5])
title20 = slide20.shapes.title
title20.text = "דופק הסנכרון (Tpw) ווריפיקציה אוטומטית"
set_rtl(title20.text_frame.paragraphs[0])
title20.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text20 = slide20.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(1.5))
tf_text20 = txBox_text20.text_frame

add_bullet(tf_text20, "• אות ה-HSYNC: זהו הטריגר הפיזי המורה למסך לבצע Retrace לתחילת השורה הבאה.")
add_bullet(tf_text20, "• דיוק של מחזור בודד: הצבת סמנים על הדילוג של האות הפיזי מוכיחה רוחב דופק מדויק של 3.84us, השקול ל-96 מחזורי שעון.")
add_bullet(tf_text20, "• Self-Checking Testbench: מעבר לבדיקה הוויזואלית, סביבת הבדיקה (TB) כוללת מנגנון אוטומטי (Assert) שבודק את רוחב הדופק בזמן ריצה ומתריע על כל חריגה מהתקן.")

img_path20 = "docs/presentation/assets/modelsim_sync_pulse_assertion.png" 
if os.path.exists(img_path20):
    slide20.shapes.add_picture(img_path20, Inches(1.0), Inches(3.5), width=Inches(8.0))
else:
    err_box = slide20.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(5), Inches(1))
    err_p = err_box.text_frame.paragraphs[0]; err_p.text = "[שומר מקום: modelsim_sync_pulse_assertion.png]"
    err_p.font.size = Pt(14); err_p.font.color.rgb = RGBColor(0, 112, 192); err_p.alignment = PP_ALIGN.CENTER
 
 
# ==========================================
# שקף 21: המעבר לחומרה - מגבלות הסימולציה והקופסה השחורה
# ==========================================
slide21 = prs.slides.add_slide(prs.slide_layouts[5])
title21 = slide21.shapes.title
title21.text = "המעבר לחומרה: מסימולציה ל'קופסה שחורה'"
set_rtl(title21.text_frame.paragraphs[0])
title21.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text21 = slide21.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
tf_text21 = txBox_text21.text_frame

add_bullet(tf_text21, "• העולם האידיאלי (ModelSim): עד כה בדקנו את הלוגיקה בסביבה וירטואלית. ראינו כל אות ומונה, אך ללא עיכובים פיזיים (Zero Latency).")
add_bullet(tf_text21, "• המציאות הפיזית (FPGA): לאחר סינתזה וצריבה, הלוגיקה הופכת ל'קופסה שחורה' אטומה. ישנם זמני התפשטות (Propagation delays) של האותות בתוך הסיליקון.")
add_bullet(tf_text21, "• מגבלת הדיבוג החיצוני: מבחוץ (למשל עם אוסצילוסקופ), ניתן למדוד רק את הפינים הפיזיים (VGA_R/G/B, Sync). אין לנו שום דרך לראות מתי המונה מתאפס או איזו כתובת נשלחת לזיכרון.")

# ציור קופסה שחורה שממחישה את המגבלה
black_box = slide21.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(4.0), Inches(4.0), Inches(1.5))
black_box.fill.solid()
black_box.fill.fore_color.rgb = RGBColor(40, 40, 40)
tb_black = black_box.text_frame
tb_black.text = "Artix-7 FPGA (קופסה שחורה)\nאין גישה לאותות פנימיים"
tb_black.paragraphs[0].font.bold = True
tb_black.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
tb_black.paragraphs[0].alignment = PP_ALIGN.CENTER
set_rtl(tb_black.paragraphs[0])

# פינים יוצאים (חיצים)
slide21.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.0), Inches(4.3), Inches(0.8), Inches(0.3))
slide21.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.0), Inches(4.9), Inches(0.8), Inches(0.3))
out_txt = slide21.shapes.add_textbox(Inches(7.8), Inches(4.2), Inches(1.5), Inches(1.0))
out_txt.text_frame.text = "פינים בלבד\n(RGB, Sync)"

# ==========================================
# שקף 22: הפתרון - Vivado ILA (קופסה לבנה)
# ==========================================
slide22 = prs.slides.add_slide(prs.slide_layouts[5])
title22 = slide22.shapes.title
title22.text = "דיבוג בתוך הסיליקון: Vivado ILA (קופסה לבנה)"
set_rtl(title22.text_frame.paragraphs[0])
title22.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox_text22 = slide22.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(9.5), Inches(1.5))
tf_text22 = txBox_text22.text_frame

add_bullet(tf_text22, "• הפתרון ההנדסי: הוספת Integrated Logic Analyzer (ILA) – רכיב חומרה פנימי שדוגם את האותות ומשדר אותם למחשב (בדיקת קופסה לבנה).")
add_bullet(tf_text22, "• היתרון (מציאות מול תיאוריה): מאפשר לראות את זמני ההשהיה האמיתיים (Latency) שבין שליחת בקשה לזיכרון (addrb) לקבלת הנתון (doutb).")
add_bullet(tf_text22, "• החיסרון (Trade-off): ה-ILA אינו חינמי. הוא 'גוזל' משאבי חומרה יקרים מהכרטיס (LUTs ו-BRAM) לצורך לוגיקת הדגימה ואגירת הנתונים.")

# טבלת תוכנית הניסוי - מה אנחנו בודקים
x2, y2, cx2, cy2 = Inches(0.5), Inches(3.8), Inches(9.0), Inches(2.5)
table_shape2 = slide22.shapes.add_table(4, 3, x2, y2, cx2, cy2)
table2 = table_shape2.table

table2.columns[0].width = Inches(2.0)
table2.columns[1].width = Inches(3.5)
table2.columns[2].width = Inches(3.5)

headers2 = ["יעד הבדיקה", "מגבלת ה'קופסה השחורה'", "האות שנדגום ב-ILA"]
for col_idx, header in enumerate(headers2):
    cell = table2.cell(0, col_idx)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(15)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_rtl(cell.text_frame.paragraphs[0])

row_data2 = [
    ["זרימת נתונים ו-Latency", "רואים רק את הפיקסל הסופי מתחלף ביציאה.", "addrb, doutb, pxl_data_reg (בדיקת צנרת)"],
    ["מיקום אופקי", "לא ניתן לדעת איזה פיקסל משורטט כרגע.", "display_x (מונה הפיקסלים)"],
    ["אכיפת שחור (Blanking)", "רואים מסך חשוך, ללא הקשר תזמוני.", "in_display_area_delayed"]
]

for row_idx, row in enumerate(row_data2):
    for col_idx, text in enumerate(row):
        cell = table2.cell(row_idx + 1, col_idx)
        cell.text = text
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        if col_idx == 0:
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        else:
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        set_rtl(cell.text_frame.paragraphs[0])
    
# ==========================================
# שמירת המצגת
# ==========================================
output_path = os.path.join(base_path, 'VGA_Presentation_Yossi.pptx')
prs.save(output_path)
print(f"All {len(prs.slides)} Slides generated and saved successfully at: {output_path}")